from pptx import Presentation
from pptx.util import Inches, Pt
import os

def add_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = content_lines[0]
    for line in content_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide

def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Add title manually
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(9))
    else:
        p = tf.add_paragraph()
        p.text = f"(Image missing: {image_path})"
    return slide

def build_deck(path="presentation/GridPulse_India.pptx"):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "GridPulse India"
    slide.placeholders[1].text = (
        "Forecasting India electricity demand with hierarchical LSTM models\n"
        "and Time Series Analysis techniques."
    )

    # 2. Where TSA Methods are Applied
    add_slide(
        prs,
        "Time Series Analysis (TSA) Methods Applied",
        [
            "Fourier Seasonality: Applied to capture periodic weekly & annual patterns.",
            "Global LSTM Networks: Used to learn shared temporal features across different hierarchy levels.",
            "Hierarchical Reconciliation (sktime): Applied BottomUp to ensure forecasts add up seamlessly from city to national levels.",
            "Time Features: Leveraging day-of-week, month, and trend signals for structured forecasting."
        ],
    )

    # 3. System Architecture Image Slide
    add_image_slide(prs, "System Architecture", "assets/system_architecture.png")

    # 4. Data Layer & Feature Engineering
    add_slide(
        prs,
        "Data & Feature Engineering Layer",
        [
            "Data Generation: Synthetic electricity demand simulated for India, Maharashtra, and Mumbai.",
            "Hierarchy: Organized structured data (National -> State -> City).",
            "Fourier Transforms: Extracting cyclic patterns (e.g. daily/weekly load fluctuations).",
            "Calendar Elements: Incorporating weekends, trends, and seasonal peaks."
        ],
    )

    # 5. Model Layer
    add_slide(
        prs,
        "Modeling & Neural Network Layer",
        [
            "Global LSTM Model: A single unified Long Short-Term Memory model trained across all regions.",
            "Multivariate Inputs: Feeding scaled Fourier and temporal features.",
            "Generalization: Learns correlations that apply globally but adapt locally.",
            "Evaluation: Monitored through Mean Absolute Error (MAE) and MSE."
        ],
    )

    # 6. Hierarchy & Reconciliation
    add_slide(
        prs,
        "Delivery & Reconciliation Layer",
        [
            "sktime Integration: Utilizing Reconciler to adjust independent forecasts.",
            "Aggregation Logic: BottomUp approach (Mumbai + other local → Maharashtra → India).",
            "Coherent Outputs: Ensuring sum of parts equals the whole.",
            "Delivery: Structured CSV forecast datasets and visual performance summaries."
        ],
    )

    prs.save(path)
    print(f"Presentation saved to {path}")

if __name__ == "__main__":
    build_deck()
